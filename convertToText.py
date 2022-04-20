import sys
from os import listdir
from os.path import isfile, join

import re

from pptx import Presentation
from docx import Document
import pdfplumber

def getPowerPointText(path):
    # list of slides as text
    slides = []
    prs = Presentation(path)
#    print(path)
#    print("----------------------")
    for slide in prs.slides:
        slideContent = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slideContent.append(shape.text)
        slides.append("\n".join(slideContent))
    return slides

def getWordText(path):
    # list of paragraphs
    doc = Document(path)
    paragraphs = []
    for para in doc.paragraphs:
        paragraphs.append(para.text)
    return paragraphs

def getPDFText(path):
    pdf = pdfplumber.open(path)
    pages = []
    for singlePage in pdf.pages:
        pages.append(singlePage.extract_text())
    return pages

onlyFiles = sorted([f for f in listdir(sys.argv[1]) if isfile(join(sys.argv[1], f))], key=str.lower)

rPPTX = re.compile(".*\.pptx$", re.IGNORECASE)
rDOCX = re.compile(".*\.docx$", re.IGNORECASE)
#rXLSX = re.compile(".*\.xlsx$", re.IGNORECASE)
rPDF = re.compile(".*\.pdf$", re.IGNORECASE)

tagList = []
for eachFile in onlyFiles:
    print(eachFile)
#    tags = [eachfile]
    onlyText = []
    if re.match(rPPTX, eachFile): onlyText = getPowerPointText(join(sys.argv[1], eachFile))
    if re.match(rDOCX, eachFile): onlyText = getWordText(join(sys.argv[1], eachFile))
#    if re.match(rXLSX, eachFile): tags = getExcelText(join(sys.argv[1], eachFile))
    if re.match(rPDF,  eachFile): onlyText = getPDFText(join(sys.argv[1], eachFile))
    if onlyText: onlyText = list(filter(None, onlyText))
    print("list length: ", len(onlyText))
#    print(onlyText)
    print('')

#    tagList.append(tags)


# for x in onlyWord:
#   print(x)
#
# for x in onlyExcel:
#   print(x)
#
# for x in onlyPDF:
#   print(x)
