import sys
import os
import re
import string
from enum import Enum
import unicodedata, itertools

import pptx
from docx import Document
import pdfplumber
import xlsxwriter
from polyglot.detect import Detector


def getPowerPointText(path):
    # list of slides as text
    slides = []
    title = ""
    year = ""
    try:
        prs = pptx.Presentation(path)
    except:
        print("--------- Error with: ", path)
        return slides, title, year
    for slide in prs.slides:
        slideContent = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slideContent.append(control_char_re.sub(' ', shape.text))
        slides.append(" ".join(slideContent))
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            if len(shape.text) > 0:
                try:
                    foundText = re.match('^([^\r\n]+)', shape.text).group(1)
                except AttributeError:
                    continue
                if foundText != "The European Patent Office":
                    title = control_char_re.sub(' ', foundText)
                    title = (title[:250] + "..") if len(title) > 250 else title
                    break
    yearRegex = re.compile('(201[4-9]|202[0-2])', re.IGNORECASE)
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            year = yearRegex.search(shape.text)
            if year:
                year = year.group(1)
                break
    return slides, title, year, len(prs.slides) + 1

def getWordText(path):
    # list of paragraphs
    paragraphs = []
    title = ""
    year = ""
    try:
        doc = Document(path)
    except:
        print("--------- Error with: ", path)
        return paragraphs, title, year
    for para in doc.paragraphs:
        paragraphs.append(para.text)
    title = control_char_re.sub(' ', doc.paragraphs[0].text)
    yearRegex = re.compile('(201[4-9]|202[0-2])', re.IGNORECASE)
    for parNum in range(min(20, len(doc.paragraphs) - 1)):
        year = yearRegex.search(doc.paragraphs[parNum].text)
        if year:
            year = year.group(1)
            break
    return paragraphs, title, year

def getPDFText(path):
    pdf = pdfplumber.open(path)
    pages = []
    for singlePage in pdf.pages:
        pages.append(singlePage.extract_text())
    year = ""
    year = re.search('(201[4-9]|202[0-2])', pdf.pages[0].extract_text())
    if year: year = year.group(1)
    return pages, year, len(pdf.pages) + 1

def get_TypeOfMaterial(filetype, onlyText):
    query1 = re.compile('exercises?', re.IGNORECASE)
    query2 = re.compile('exercise.{,5}solution', re.IGNORECASE)
    query3 = re.compile('case stud', re.IGNORECASE)
    query4 = re.compile('cheat sheet', re.IGNORECASE)
    query5 = re.compile('mock', re.IGNORECASE)
    result = []
    if filetype == extension.PPTX: result.append("Presentation|2dc089fb-8444-4cd3-a9ca-9fcc728aac7a; ")
    if any(query1.search(line) for line in onlyText): result.append("Exercise|cac0e695-c696-42d9-acac-27fc80b91501; ")
    if any(query2.search(line) for line in onlyText): result.append("Exercise - solutions|68e2a42d-6972-42be-bc70-6af1335b444c; ")
    if any(query3.search(line) for line in onlyText): result.append("Case study|4e679404-4a79-4ace-9bdb-d5f72cd66b0d; ")
    if any(query4.search(line) for line in onlyText): result.append("Cheat sheet|833de978-768a-4a78-a2fa-1f4bebe0d6a7; ")
    if any(query5.search(line) for line in onlyText): result.append("Mock|eba4b3c5-2cd7-41bc-b7fd-487d6d5586d5; ")
    return ''.join(result)

def get_ip_strategy(onlyText):
    query = re.compile('(?:ip|intellectual property).{1,4}strat?g??', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_patent_licencing(onlyText):
    query = re.compile('li[cs]en[csz]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_patent_landscape(onlyText):
    query = re.compile('landscap', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_patent_valuation(onlyText):
    query = re.compile('(?:valuation|bewert|valoris)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_boa_decisions(onlyText):
    query = re.compile('[gtj][, ][0-9][0-9][0-9][0-9]/[0-9][0-9]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_classification(onlyText):
    query = re.compile('classif', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_amendments(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}(?:123|76)|(?:richtlinie|guideline|directive|GL) ?H)', re.IGNORECASE)
    query2 = re.compile('amendment', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).{0,10}(?:19\(2\)|34\(2.?b\)).?PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_clarity(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}84|F-IV.{0,3}4|(?:richtlinie|guideline|directive|GL) ?F[- ]?IV.{0,3}4)', re.IGNORECASE)
    query2 = re.compile('(?:clarity|suff.+ (of)? disclos|broad claim|lack of support|concise)', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).{0,10}6.?PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_exclusions(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}53|(?:r|rule|regle|regel).{0,10}(?:28|29)|(?:richtlinie|guideline|directive|GL) ?G[- ]?II.{0,3}[345])', re.IGNORECASE)
    query2 = re.compile('(?:(?:exception|exclusion).{1,5}patentability|ordre public|treatments?(?:\s\S){0,4} bod|(?:surgery|therapy|diagnostic)(?:\s\S){0,4} (?:human|animal)|human embryo|clon.{1,3}(?:\s\S){0,2} human|human(?:\s\S){0,4} clon.{1,3})', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).(?:9\.1|39\.1|67\.1).{0,8}PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_inventiveness(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}56|(?:richtlinie|guideline|directive|GL) ?G[- ]?VII[^IV])', re.IGNORECASE)
    query2 = re.compile('(?:inventive step|inventiveness|erfinderische Tätigkeit|activité inventive)', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).{0,10}33.3.{0,8}PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_novelty(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}54|F-IV|(?:richtlinie|guideline|directive|GL) ?F[- ]?IV)', re.IGNORECASE)
    query2 = re.compile('(?:novelty|nouveauté|neuheit)', re.IGNORECASE)
    query3 = re.compile('(?:R|Rule).{0,10}33\.1.?[abc].{0,10}PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_priority(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}(?:87|88|89)|F-VI[^IV]|(?:richtlinie|guideline|directive|GL) ?F[- ]?VI[^IV])', re.IGNORECASE)
    query2 = re.compile('(?:priority|priorität|priorité|state of the art|stand der technik|art antérieur)', re.IGNORECASE)
    query3 = re.compile('(?:(?:A|Art).{0,10}8 ?PCT|(?:R|Rule).{0,10}64\.1.?PCT)', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_sufficiency(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}83|(?:richtlinie|guideline|directive|GL) ?F[- ]?III)', re.IGNORECASE)
    query2 = re.compile('(?:F-III|sufficien.{1,3}(?:\s\S){0,3} disclos|offenbarung)', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).{0,10}5 ?PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_unity(onlyText):
    query1 = re.compile('(?:(?:A|Art).{0,10}82|(?:richtlinie|guideline|directive|GL) ?F[- ]?V[^IV]|F-V[^IV])', re.IGNORECASE)
    query2 = re.compile('(?:G ?2/92|unity of invention|require.{1,5}(?:\s\S){0,2} unity|einheitlichkeit|unité)', re.IGNORECASE)
    query3 = re.compile('(?:A|Art).{0,10}13\.1 ?PCT', re.IGNORECASE)
    epc = 0
    pct = 0
    common = 0
    for i, item in enumerate(onlyText):
        if query1.search(item):
            epc += 1
            common += 1
        if query2.search(item): common += 1
        if query3.search(item):
            pct += 1
            common += 1
    return [common, epc, pct]

def get_pct(onlyText):
    query = re.compile('(?:chapter I|chapter II|PCT procedure)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_national(onlyText):
    query = re.compile('(?:national|NATL|USPTO|CIPO|SIPO|KIPO|JPO|INPI|DPMA)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_litigation(onlyText):
    query = re.compile('(?:litig+|infring|contre.?fa)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_enforcement(onlyText):
    query = re.compile('enforcement?', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_revocation(onlyText):
    query = re.compile('r.?vo.?at', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_eqe(onlyText):
    query = re.compile('(?:eqe|apec)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_diversity(onlyText):
    query = re.compile('(?:D&I|diversity(?:-| | and )?incl)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_data_protection(onlyText):
    query = re.compile('(?:DPO|GDPR|data protection)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def get_wellbeing(onlyText):
    query = re.compile('(?:well.?be|mental.?health)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return len(matchesIndex)

def writeFields(fields, worksheetOut):
    column = 0
    for field in fields:
        worksheetOut.write(row, column, field)
        column += 1



topics = [
    "IP strategy|64267029-dfa1-4a94-b27f-67a407103c8c",
    "Patent licencing, transfer and dissemination|b0462b20-1e15-4310-a472-575b60375532",
    "Patent landscape analysis|1e6fcc69-7415-4fab-9e53-e2076a6878af",
    "Patent valuation|704f1086-252d-44a2-ba6a-069c388694c2",
    "Board of Appeal decisions|9fc885a3-b9cd-4a7d-a31b-22617def861e",
    "Classification and documentation|9f130a00-885e-4cf1-989e-6ac68658eafd",
    "Amendments|6c67b21e-8654-4e5a-9c47-f8a6c0251aa1",
    "Clarity|2c950022-e6b1-4a4b-b788-d180e4f156ac",
    "Exclusions to patentability|17cb3b3f-441b-481b-8317-88f6e7966d07",
    "Inventiveness|d26816ab-a81c-4d89-b4eb-cee1abab59f7",
    "Novelty|595bb310-797c-4b47-b137-3a36c39be08c",
    "Priority|bb3af4c3-0efe-4d4e-8bfb-33c5df715ae3",
    "Sufficiency of disclosure|dc2aefb9-8d77-415b-8be2-732dd66e276e",
    "Unity of invention|1d5222d0-9873-4333-93cf-712be08fe02a",
    "PCT specific procedures|08e842f3-de6c-4b8c-a844-c2bbcd8c6369",
    "National specific procedures|fd9c1e21-8ce8-41f5-845c-bb2688396b42",
    "Litigation and infringement procedures in Europe|cd7ae951-29ca-4c65-a8ad-98c0bb61a638",
    "Enforcement measures|19531776-8f1b-4ba6-96ec-cda92be380dc",
    "Validity topics in National revocation procedure|0a913a75-a0ce-4980-befa-b4b5a0f5bcaf",
    "Software for Certifications (EQE/APEC…)|8de9260b-9c0e-47ee-9cea-f6f806fa2300",
    "Diversity and Inclusion|c921bbf7-d0d6-4500-b519-a3c2514571dd",
    "Data Protection|6266e964-0e74-46b3-bd74-0c54995ca7de",
    "Mental Well-being|b41b1481-41ff-4613-a353-ad9e38d06f6e"
    ]

class extension(Enum):
    PPTX = 1
    DOCX = 2
    PDF  = 3

thresholdMatches = 3

onlyFiles = sorted([f for f in os.listdir(sys.argv[1]) if os.path.isfile(os.path.join(sys.argv[1], f))], key=str.lower)
outputFile = re.sub('^[./]+', '', sys.argv[1])
outputFile = re.sub('^files/', '', outputFile)
outputFile = re.sub('/$', '', outputFile)
outputFile = re.sub('/', '-', outputFile)

workbook = xlsxwriter.Workbook('outputs/' + outputFile + '.xlsx')
worksheet = workbook.add_worksheet()
row = 0
column = 0
content = [
    "File name", "Pages", "Title", "Type of material",
    "Language", "Year", "Topic",
    "Patent law ref", "ip_strategy", "patent_licencing",
    "patent_landscape", "patent_valuation", "boa_decisions",
    "classification", "amendments", "clarity",
    "exclusions", "inventiveness", "novelty",
    "priority", "sufficiency", "unity", "PCT",
    "national", "litigation", "enforcement",
    "revocation", "eqe", "diversity",
    "data_protection", "wellbeing"]
for item in content :
    worksheet.write(row, column, item)
    column += 1
row += 1

rPPTX = re.compile(".*\.pptx$", re.IGNORECASE)
rDOCX = re.compile(".*\.docx$", re.IGNORECASE)
rPDF = re.compile(".*\.pdf$", re.IGNORECASE)

all_chars = (chr(i) for i in range(sys.maxunicode))
categories = {'Cc'}
control_chars = ''.join(map(chr, itertools.chain(range(0x00,0x20), range(0x7f,0xa0))))
control_char_re = re.compile('[%s]' % re.escape(control_chars))

usefulFiles = 0
for eachFile in onlyFiles:
#    print(eachFile)
    fields = [""] * 31
    fullPath = os.path.join(sys.argv[1], eachFile)
    fields[0] = eachFile

    onlyText = []
    filetype = ""
    if rPPTX.match(eachFile):
        onlyText, fields[2], fields[5], fields[1] = getPowerPointText(fullPath)
        filetype = extension.PPTX
    if rDOCX.match(eachFile):
        onlyText, fields[2], fields[5] = getWordText(fullPath)
        filetype = extension.DOCX
    if rPDF.match(eachFile):
        onlyText, fields[5], fields[1] = getPDFText(fullPath)
        filetype = extension.PDF

    if onlyText: onlyText = list(filter(None, onlyText))
    if not len(onlyText):
        writeFields(fields, worksheet)
        row += 1
        continue
    usefulFiles += 1

    fields[3] = get_TypeOfMaterial(filetype, onlyText)
    fields[4] = Detector(''.join(x for x in "".join(onlyText[0:min(50,len(onlyText))]) if x.isprintable()), quiet=True).language.code.upper()

    if len(onlyText) > 150:
        writeFields(fields, worksheet)
        row += 1
        continue

    fields[8] = get_ip_strategy(onlyText)
    fields[9] = get_patent_licencing(onlyText)
    fields[10] = get_patent_landscape(onlyText)
    fields[11] = get_patent_valuation(onlyText)
    fields[12] = get_boa_decisions(onlyText)
    fields[13] = get_classification(onlyText)

    epcMatches = 0
    pctMatches = 0

    amendments = get_amendments(onlyText)
    fields[14] = amendments[0]
    epcMatches += amendments[1]
    pctMatches += amendments[2]
    if amendments[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 123 EPC; Art. 76 EPC; "
    if amendments[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 19(2) PCT, Art. 34(2)(b) PCT; "

    clarity = get_clarity(onlyText)
    fields[15] = clarity[0]
    epcMatches += clarity[1]
    pctMatches += clarity[2]
    if clarity[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 84 EPC; "
    if clarity[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 6 PCT; "

    exclusions = get_exclusions(onlyText)
    fields[16] = exclusions[0]
    epcMatches += exclusions[1]
    pctMatches += exclusions[2]
    if exclusions[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 53 EPC; Rule 28 EPC; Rule 29 EPC; "
    if exclusions[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 9.1 PCT; Rule 39.1 PCT; Rule 67.1 PCT; "

    inventiveness = get_inventiveness(onlyText)
    fields[17] = inventiveness[0]
    epcMatches += inventiveness[1]
    pctMatches += inventiveness[2]
    if inventiveness[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 56 EPC; "
    if inventiveness[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 33(3) PCT; "

    novelty = get_novelty(onlyText)
    fields[18] = novelty[0]
    epcMatches += novelty[1]
    pctMatches += novelty[2]
    if novelty[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 54 EPC; "
    if novelty[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 33.1 PCT; "

    priority = get_priority(onlyText)
    fields[19] = priority[0]
    epcMatches += priority[1]
    pctMatches += priority[2]
    if priority[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 87 EPC; Art. 88 EPC; Art. 89 EPC; "
    if priority[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 8 PCT; Rule 64.1 PCT; "

    sufficiency = get_sufficiency(onlyText)
    fields[20] = sufficiency[0]
    epcMatches += sufficiency[1]
    pctMatches += sufficiency[2]
    if sufficiency[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 83 EPC; "
    if sufficiency[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 5 PCT; "

    unity = get_unity(onlyText)
    fields[21] = unity[0]
    epcMatches += unity[1]
    pctMatches += unity[2]
    if unity[1] > thresholdMatches:
        fields[7] = fields[7] + "Art. 82 EPC; "
    if unity[2] > thresholdMatches:
        fields[7] = fields[7] + "Art. 13.1 PCT; "

    pct = get_pct(onlyText)
    pctMatches += pct
    fields[22] = pctMatches

    fields[23] = get_national(onlyText)
    fields[24] = get_litigation(onlyText)
    fields[25] = get_enforcement(onlyText)
    fields[26] = get_revocation(onlyText)
    fields[27] = get_eqe(onlyText)
    fields[28] = get_diversity(onlyText)
    fields[29] = get_data_protection(onlyText)
    fields[30] = get_wellbeing(onlyText)

    maxMatches = max(fields[8:30])
#    print(maxMatches)
    maxMatchesIndexes = [i for i, j in enumerate(fields[8:30]) if j == maxMatches]
    if len(maxMatchesIndexes) > 1:
        fields[6] = " "
    else:
        # assigns the topic to the concept with most matches
        fields[6] = topics[maxMatchesIndexes[0]]
        # if the topic is in the patent law concepts check the title
        if maxMatchesIndexes[0] > 5 and maxMatchesIndexes[0] < 14:
            titleMatches = [''] * 8
            titleMatches[0] = get_amendments([fields[2]])[0]
            titleMatches[1] = get_clarity([fields[2]])[0]
            titleMatches[2] = get_exclusions([fields[2]])[0]
            titleMatches[3] = get_inventiveness([fields[2]])[0]
            titleMatches[4] = get_novelty([fields[2]])[0]
            titleMatches[5] = get_priority([fields[2]])[0]
            titleMatches[6] = get_sufficiency([fields[2]])[0]
            titleMatches[7] = get_unity([fields[2]])[0]
            # check if there is only one matching concept in the title
            nonNullMatches = [i for i, j in enumerate(titleMatches) if j > 0]
#            print(nonNullMatches)
            # if only one found, assign the topic to that one
            if len(nonNullMatches) == 1:
                fields[6] = topics[6 + nonNullMatches[0]]
            # if there are multiple matching concepts, check if the fulltext matches one more than 3x the other ones
            # we know already from the checks at the beginning that only one field has the max value
            if len(nonNullMatches) > 1 or len(nonNullMatches) == 0:
                for i, j in enumerate(fields[14:21]):
            # if any value is higher than 1/3 the max, assign generic topic
                    if j > maxMatches/3 and i != (maxMatchesIndexes[0] - 6):
                        fields[6] = "Patent law concepts|493f6ca1-16fd-4f96-bcd7-e46f81984678"

    column = 0
    for field in fields:
#        print(row, " ", column, " ", field)
        worksheet.write(row, column, field)
        column += 1

    row += 1

workbook.close()
if not usefulFiles : os.remove('outputs/' + outputFile + '.xlsx')
