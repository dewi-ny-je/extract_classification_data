import sys
from os import listdir
from os.path import isfile, join

import re
import string
from enum import Enum

from pptx import Presentation
from docx import Document
import pdfplumber
import xlsxwriter

from polyglot.detect import Detector


def getPowerPointText(path):
    # list of slides as text
    slides = []
    title = ""
    prs = Presentation(path)
#    print(path)
#    print("----------------------")
    for slide in prs.slides:
        slideContent = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slideContent.append(shape.text)
        slides.append("\n".join(slideContent))
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text"):
            title = shape.text
            break
    return slides, title

def getWordText(path):
    # list of paragraphs
    doc = Document(path)
    paragraphs = []
    title = ""
    for para in doc.paragraphs:
        paragraphs.append(para.text)
    title = doc.paragraphs[0].text
    return paragraphs, title

def getPDFText(path):
    pdf = pdfplumber.open(path)
    pages = []
    for singlePage in pdf.pages:
        pages.append(singlePage.extract_text())
    return pages

def get_TypeOfMaterial(filetype, onlyText):
    query1 = re.compile('exercises', re.IGNORECASE)
    query2 = re.compile('exercise.{,5}solution', re.IGNORECASE)
    query3 = re.compile('case stud', re.IGNORECASE)
    query4 = re.compile('cheat sheet', re.IGNORECASE)
    query5 = re.compile('mock', re.IGNORECASE)
    result = []
    if filetype == extension.PPTX: result.append("Presentation")
    if any(query1.search(line) for line in onlyText): result.append("Exercise")
    if any(query2.search(line) for line in onlyText): result.append("Exercise - solutions")
    if any(query3.search(line) for line in onlyText): result.append("Case study")
    if any(query4.search(line) for line in onlyText): result.append("Cheat sheet")
    if any(query5.search(line) for line in onlyText): result.append("Mock")
    return ', '.join(result)

def get_ip_strategy(onlyText):
    query = re.compile('(?:ip|intellectual property).{1,4}strat?g??', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_patent_licencing(onlyText):
    query = re.compile('li[cs]en[csz]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_patent_landscape(onlyText):
    query = re.compile('landscap', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_patent_valuation(onlyText):
    query = re.compile('(?:valuation|bewert|valoris)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_boa_decisions(onlyText):
    query = re.compile('[gtj][, ][0-9][0-9][0-9][0-9]/[0-9][0-9]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_classification(onlyText):
    query = re.compile('classif', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_amendments(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}(?:123|76)', re.IGNORECASE)
    query2 = re.compile('amendment', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?H', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_clarity(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}84', re.IGNORECASE)
    query2 = re.compile('(?:clarity|suff.+ (of)? disclos|broad claim|lack of support|concise)', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?F[- ]?IV.{0,3}4', re.IGNORECASE)
    query4 = re.compile('F-IV.{0,3}4', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item) or query4.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_exclusions(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}53', re.IGNORECASE)
    query2 = re.compile('(?:r|rule|regle|regel).{0,10}(?:28|29)', re.IGNORECASE)
    query3 = re.compile('(?:exception|exclusion).{1,5}patentability', re.IGNORECASE)
    query4 = re.compile('(?:ordre public|treatments?(?:\s\S){0,4} bod|(?:surgery|therapy|diagnostic)(?:\s\S){0,4} (?:human|animal)|human embryo|clon.{1,3}(?:\s\S){0,2} human|human(?:\s\S){0,4} clon.{1,3})', re.IGNORECASE)
    query5 = re.compile('(?:richtlinie|guideline|directive|GL) ?G[- ]?II.{0,3}[345]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item) or query4.search(item) or query5.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_inventiveness(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}56', re.IGNORECASE)
    query2 = re.compile('(?:inventive step|inventiveness|erfinderische Tätigkeit|activité inventive)', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?G[- ]?VII[^IV]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_novelty(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}54', re.IGNORECASE)
    query2 = re.compile('(?:novelty|nouveauté|neuheit|F-IV)', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?F[- ]?IV', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_priority(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}(?:87|88|89)', re.IGNORECASE)
    query2 = re.compile('(?:priority|priorität|priorité|F-VI[^IV]|state of the art|stand der technik|art antérieur)', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?F[- ]?VI[^IV]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_sufficiency(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}83', re.IGNORECASE)
    query2 = re.compile('(?:F-III|sufficien.{1,3}(?:\s\S){0,3} disclos|offenbarung)', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?F[- ]?III', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_unity(onlyText):
    query1 = re.compile('(?:A|Art).{0,10}82', re.IGNORECASE)
    query2 = re.compile('(?:G ?2/92|unity of invention|require.{1,5}(?:\s\S){0,2} unity|einheitlichkeit|unité|F-V[^IV])', re.IGNORECASE)
    query3 = re.compile('(?:richtlinie|guideline|directive|GL) ?F[- ]?V[^IV]', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query1.search(item) or query2.search(item) or query3.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_national(onlyText):
    query = re.compile('(?:national|NATL|USPTO|CIPO|SIPO|KIPO|JPO|INPI|DPMA)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_litigation(onlyText):
    query = re.compile('(?:litig+|infring|contre.?fa)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_enforcement(onlyText):
    query = re.compile('enforcement?', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_revocation(onlyText):
    query = re.compile('r.?vo.?at', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_eqe(onlyText):
    query = re.compile('(?:eqe|apec)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_diversity(onlyText):
    query = re.compile('(?:D&I|diversity(?:-| | and )?incl)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_data_protection(onlyText):
    query = re.compile('(?:DPO|GDPR|data protection)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]

def get_wellbeing(onlyText):
    query = re.compile('(?:well.?be|mental.?health)', re.IGNORECASE)
    matchesIndex = [i for i, item in enumerate(onlyText) if query.search(item)]
    return [str(len(matchesIndex)), ", ".join(map(str, matchesIndex))]



class extension(Enum):
    PPTX = 1
    DOCX = 2
    PDF  = 3

onlyFiles = sorted([f for f in listdir(sys.argv[1]) if isfile(join(sys.argv[1], f))], key=str.lower)
outputFile = re.sub('^[./]+', '', sys.argv[1])
outputFile = re.sub('^files/', '', outputFile)
outputFile = re.sub('/$', '', outputFile)
outputFile = re.sub('/', '-', outputFile)

workbook = xlsxwriter.Workbook(outputFile + '.xlsx')
worksheet = workbook.add_worksheet()
row = 0
column = 0
content = ["File name", "path", "Language", "Title", "Type of material",
    "ip_strategy count", "ip_strategy list", "patent_licencing count",
    "patent_licencing list", "patent_landscape count",
    "patent_landscape list", "patent_valuation count", "patent_valuation list",
    "boa_decisions count", "boa_decisions list", "classification count",
    "classification list", "amendments count", "amendments list",
    "clarity count", "clarity list", "exclusions count", "exclusions list",
    "inventiveness count", "inventiveness list", "novelty count",
    "novelty list", "priority count", "priority list", "sufficiency count",
    "sufficiency list", "unity count", "unity list", "national count",
    "national list", "litigation count", "litigation list",
    "enforcement count", "enforcement list", "revocation count",
    "revocation list", "eqe count", "eqe list", "diversity count",
    "diversity list", "data_protection count", "data_protection list",
    "wellbeing count", "wellbeing list"]

# iterating through content list
for item in content :

    # write operation perform
    worksheet.write(row, column, item)

    # incrementing the value of row by one
    # with each iterations.
    column += 1
row += 1

rPPTX = re.compile(".*\.pptx$", re.IGNORECASE)
rDOCX = re.compile(".*\.docx$", re.IGNORECASE)
rPDF = re.compile(".*\.pdf$", re.IGNORECASE)

#tagList = []
for eachFile in onlyFiles:
    print(eachFile)
    fullPath = join(sys.argv[1], eachFile)
    worksheet.write(row, 0, eachFile)
    worksheet.write(row, 1, fullPath)

    onlyText = []
    title = ""
    filetype = ''
    if re.match(rPPTX, eachFile):
        onlyText, title = getPowerPointText(fullPath)
        filetype = extension.PPTX
    if re.match(rDOCX, eachFile):
        onlyText, title = getWordText(fullPath)
        filetype = extension.DOCX
    if re.match(rPDF,  eachFile):
        onlyText = getPDFText(fullPath)
        filetype = extension.PDF

    if onlyText: onlyText = list(filter(None, onlyText))
    if not len(onlyText): continue

    worksheet.write(row, 2, Detector(''.join(x for x in "".join(onlyText) if x.isprintable())).language.code.upper())
    worksheet.write(row, 3, title)
    worksheet.write(row, 4, get_TypeOfMaterial(filetype, onlyText))

    output = get_ip_strategy(onlyText)
    worksheet.write(row, 5, output[0])
    worksheet.write(row, 6, output[1])

    output = get_patent_licencing(onlyText)
    worksheet.write(row, 7, str(output[0]))
    worksheet.write(row, 8, output[1])

    output = get_patent_landscape(onlyText)
    worksheet.write(row, 9, output[0])
    worksheet.write(row, 10, output[1])

    output = get_patent_valuation(onlyText)
    worksheet.write(row, 11, output[0])
    worksheet.write(row, 12, output[1])

    output = get_boa_decisions(onlyText)
    worksheet.write(row, 13, output[0])
    worksheet.write(row, 14, output[1])

    output = get_classification(onlyText)
    worksheet.write(row, 15, output[0])
    worksheet.write(row, 16, output[1])

    output = get_amendments(onlyText)
    worksheet.write(row, 17, output[0])
    worksheet.write(row, 18, output[1])

    output = get_clarity(onlyText)
    worksheet.write(row, 19, output[0])
    worksheet.write(row, 20, output[1])

    output = get_exclusions(onlyText)
    worksheet.write(row, 21, output[0])
    worksheet.write(row, 22, output[1])

    output = get_inventiveness(onlyText)
    worksheet.write(row, 23, output[0])
    worksheet.write(row, 24, output[1])

    output = get_novelty(onlyText)
    worksheet.write(row, 25, output[0])
    worksheet.write(row, 26, output[1])

    output = get_priority(onlyText)
    worksheet.write(row, 27, output[0])
    worksheet.write(row, 28, output[1])

    output = get_sufficiency(onlyText)
    worksheet.write(row, 29, output[0])
    worksheet.write(row, 30, output[1])

    output = get_unity(onlyText)
    worksheet.write(row, 31, output[0])
    worksheet.write(row, 32, output[1])

    output = get_national(onlyText)
    worksheet.write(row, 33, output[0])
    worksheet.write(row, 34, output[1])

    output = get_litigation(onlyText)
    worksheet.write(row, 35, output[0])
    worksheet.write(row, 36, output[1])

    output = get_enforcement(onlyText)
    worksheet.write(row, 37, output[0])
    worksheet.write(row, 38, output[1])

    output = get_revocation(onlyText)
    worksheet.write(row, 39, output[0])
    worksheet.write(row, 40, output[1])

    output = get_eqe(onlyText)
    worksheet.write(row, 41, output[0])
    worksheet.write(row, 42, output[1])

    output = get_diversity(onlyText)
    worksheet.write(row, 43, output[0])
    worksheet.write(row, 44, output[1])

    output = get_data_protection(onlyText)
    worksheet.write(row, 45, output[0])
    worksheet.write(row, 46, output[1])

    output = get_wellbeing(onlyText)
    worksheet.write(row, 47, output[0])
    worksheet.write(row, 48, output[1])

    row += 1
#    print(tags)
#    tagList.append(';'.join(tags))

#print(tagList)
workbook.close()
